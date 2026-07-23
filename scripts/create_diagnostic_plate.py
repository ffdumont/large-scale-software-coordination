"""Build the coordination-diagnostic plate as a 16:9 PowerPoint deck.

Mirrors research/decision-engine-plate.html: an engineer's survey plate for the
importance x performance matrix. Two slides — the field itself, then how to read it.

Run:  uv run --with python-pptx scripts/create_diagnostic_plate.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "outputs" / "coordination-diagnostic-plate.pptx"

# ---------------------------------------------------------------- design tokens

GROUND = RGBColor(0xEF, 0xF1, 0xF3)  # drafting film, cool off-white
INK = RGBColor(0x10, 0x16, 0x1B)
INK2 = RGBColor(0x4C, 0x59, 0x64)
INK3 = RGBColor(0x7C, 0x88, 0x94)
RULE = RGBColor(0xC8, 0xD0, 0xD7)

KEEP = RGBColor(0x2F, 0x6B, 0x4F)
FIX = RGBColor(0xB0, 0x43, 0x2B)
CUT = RGBColor(0x6F, 0x7C, 0x88)
SIMPLIFY = RGBColor(0xA6, 0x77, 0x22)

# tints deepened slightly vs. the web plate — projectors wash pale fills out
KEEP_TINT = RGBColor(0xE6, 0xF0, 0xEA)
FIX_TINT = RGBColor(0xF7, 0xE7, 0xE1)
CUT_TINT = RGBColor(0xEC, 0xEF, 0xF2)
SIMPLIFY_TINT = RGBColor(0xF5, 0xEC, 0xD9)

KEEP_HATCH = RGBColor(0xD1, 0xE1, 0xD7)
FIX_HATCH = RGBColor(0xE8, 0xD1, 0xC7)

DISPLAY = "Arial Narrow"  # drafting lettering
BODY = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)

# ---------------------------------------------------------------- primitives


def track(run, hundredths_pt):
    """Letter-spacing — not exposed by python-pptx, so set a:rPr/@spc directly."""
    run.font._rPr.set("spc", str(int(hundredths_pt)))


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, vertical=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if vertical:
        tf._txBody.bodyPr.set("vert", vertical)
    return tf


def write(
    tf,
    text,
    *,
    font=BODY,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    spacing=None,
    space_before=0,
    line=None,
    first=False,
):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    if line:
        para.line_spacing = line
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing:
        track(run, spacing)
    return para


def rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0.0):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    return shape


def hairline(slide, x, y, w, color=RULE, pt=1.0):
    rect(slide, x, y, w, Pt(pt), fill=color)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GROUND
    return slide


def footer(slide, n):
    hairline(slide, M, Inches(6.98), W - 2 * M, RULE, 0.75)
    tf = textbox(slide, M, Inches(7.08), Inches(7.0), Inches(0.26))
    write(tf, "Coordination Model Diagnostic", font=MONO, size=9, color=INK3, first=True)
    tf = textbox(slide, W - M - Inches(1.2), Inches(7.08), Inches(1.2), Inches(0.26))
    write(tf, f"PLATE {n}/2", font=MONO, size=9, color=INK3, align=PP_ALIGN.RIGHT, first=True)


def head(slide, title, lede, designation):
    tf = textbox(slide, M, Inches(0.36), Inches(8.9), Inches(0.70))
    write(tf, title, font=DISPLAY, size=34, bold=True, color=INK, spacing=20, first=True)

    tf = textbox(slide, M, Inches(1.06), Inches(9.4), Inches(0.52))
    write(tf, lede, size=13, color=INK2, line=1.28, first=True)

    tf = textbox(slide, W - M - Inches(3.5), Inches(0.44), Inches(3.5), Inches(0.6))
    for i, part in enumerate(designation):
        write(
            tf,
            part,
            font=MONO,
            size=9,
            color=INK3,
            align=PP_ALIGN.RIGHT,
            spacing=40,
            line=1.5,
            first=(i == 0),
        )

    hairline(slide, M, Inches(1.70), W - 2 * M, INK, 1.75)


# ---------------------------------------------------------------- slide 1


QUADRANTS = [
    # col, row, coord, verdict, gloss lead, gloss emphasis, colors
    # gloss is split across two lines on purpose: the emphasis reads as its own
    # instruction and the four boxes stay typographically aligned
    (0, 0, "Important · broken", "FIX", "Load-bearing, and it is failing.",
     "Protect the function, change the form.", FIX, FIX_TINT, FIX_HATCH),
    (1, 0, "Important · working", "KEEP", "Load-bearing and holding.",
     "Do not touch it.", KEEP, KEEP_TINT, KEEP_HATCH),
    (0, 1, "Not important · broken", "CUT", "Overhead nobody relies on.",
     "The safe early win.", CUT, CUT_TINT, None),
    (1, 1, "Not important · working", "SIMPLIFY",
     "Over-served — we pay more than it is worth.",
     "Trim the effort.", SIMPLIFY, SIMPLIFY_TINT, None),
]

FIELD_X, FIELD_Y = Inches(1.16), Inches(1.90)
FIELD_W, FIELD_H = W - M - FIELD_X, Inches(4.20)
COL_W, ROW_H = Emu(int(FIELD_W / 2)), Emu(int(FIELD_H / 2))


def axis_label(slide, x, y, w, h, text, *, size, color, spacing, vertical=None, align=PP_ALIGN.CENTER, font=MONO, bold=False):
    tf = textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, vertical=vertical)
    write(tf, text, font=font, size=size, bold=bold, color=color, align=align, spacing=spacing, first=True)


def build_field(slide):
    for col, row, coord, verdict, lead, emph, color, tint, hatch in QUADRANTS:
        x = Emu(FIELD_X + col * COL_W)
        y = Emu(FIELD_Y + row * ROW_H)
        cell = rect(slide, x, y, COL_W, ROW_H, fill=tint, line_color=INK, line_pt=0.75)

        # top row = load-bearing: hatched, as "material to remain" on a survey drawing
        if hatch is not None:
            cell.fill.patterned()
            cell.fill.pattern = MSO_PATTERN.LIGHT_UPWARD_DIAGONAL
            cell.fill.fore_color.rgb = hatch
            cell.fill.back_color.rgb = tint

        pad = Inches(0.28)
        inner_w = Emu(COL_W - 2 * pad)

        tf = textbox(slide, Emu(x + pad), Emu(y + Inches(0.20)), inner_w, Inches(0.22))
        write(tf, coord.upper(), font=MONO, size=9, color=INK3, spacing=60, first=True)

        tf = textbox(slide, Emu(x + pad), Emu(y + Inches(0.48)), inner_w, Inches(0.72))
        write(tf, verdict, font=DISPLAY, size=44, bold=True, color=color, spacing=25, first=True)

        tf = textbox(slide, Emu(x + pad), Emu(y + Inches(1.28)), inner_w, Inches(0.62))
        for j, (chunk, is_emph) in enumerate(((lead, False), (emph, True))):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.line_spacing = 1.24
            run = para.add_run()
            run.text = chunk
            run.font.name = BODY
            run.font.size = Pt(13)
            run.font.bold = is_emph
            run.font.color.rgb = INK if is_emph else INK2

    rect(slide, FIELD_X, FIELD_Y, FIELD_W, FIELD_H, fill=None, line_color=INK, line_pt=2.25)


def slide_field(prs):
    slide = new_slide(prs)
    head(
        slide,
        "Two scores per function. The gap decides.",
        "Each of the twelve coordination functions is rated twice, independently — how important it is, "
        "and how well it works today. Neither number means anything alone.",
        ["COORDINATION MODEL DIAGNOSTIC", "IMPORTANCE × PERFORMANCE"],
    )

    build_field(slide)

    # y axis — reads bottom-to-top
    ax = M
    aw = Inches(0.46)
    axis_label(slide, ax, FIELD_Y, aw, Inches(1.30), "5 · ESSENTIAL", size=9, color=INK3, spacing=40, vertical="vert270")
    axis_label(slide, ax, Emu(FIELD_Y + Inches(1.42)), aw, Inches(1.40), "IMPORTANCE",
               size=13, color=INK, spacing=170, vertical="vert270", font=DISPLAY, bold=True)
    axis_label(slide, ax, Emu(FIELD_Y + FIELD_H - Inches(1.30)), aw, Inches(1.30), "1 · NICE-TO-HAVE",
               size=9, color=INK3, spacing=40, vertical="vert270")

    # x axis
    ay = Emu(FIELD_Y + FIELD_H + Inches(0.10))
    axis_label(slide, FIELD_X, ay, Inches(3.0), Inches(0.26), "1 · VERY POORLY",
               size=9, color=INK3, spacing=40, align=PP_ALIGN.LEFT)
    axis_label(slide, Emu(FIELD_X + Inches(3.2)), ay, Emu(FIELD_W - Inches(6.4)), Inches(0.26),
               "WORKS TODAY", size=13, color=INK, spacing=170, font=DISPLAY, bold=True)
    axis_label(slide, Emu(FIELD_X + FIELD_W - Inches(3.0)), ay, Inches(3.0), Inches(0.26), "5 · VERY WELL",
               size=9, color=INK3, spacing=40, align=PP_ALIGN.RIGHT)

    tf = textbox(slide, FIELD_X, Emu(FIELD_Y + FIELD_H + Inches(0.50)), FIELD_W, Inches(0.26))
    write(
        tf,
        "Hatched row = load-bearing.  Twelve functions plot onto this field once the survey closes — one point each.",
        font=MONO, size=9.5, color=INK3, align=PP_ALIGN.CENTER, spacing=20, first=True,
    )

    footer(slide, 1)
    slide.notes_slide.notes_text_frame.text = (
        "Draw the two axes in the air before this comes up — importance, and does it work today.\n\n"
        "Every function is scored twice, independently. Neither number tells you anything alone; the gap does.\n\n"
        "Important and working: keep, don't touch it. Important and broken: fix — protect the function, "
        "change the form. That is the interesting box, and it is where most of the frustration lives. "
        "Not important and broken: cut, free win. Not important but working fine: we are over-serving it, trim.\n\n"
        "The hatched top row is the load-bearing row. FIX and KEEP sit together on purpose — FIX is broken "
        "structure, not decoration. That distinction is the one that saves us in month six."
    )
    return slide


# ---------------------------------------------------------------- slide 2

NOTES = [
    (
        "COST OVERLAY",
        [
            ("A third question asks which two or three functions ", "cost more than they are worth", "."
             "  It moves points; it does not create them."),
        ],
        [
            (CUT, "Not important and costly → cut first."),
            (FIX, "Important and costly → fix it. Do not kill it."),
        ],
    ),
    (
        "PLOTTED PER ROLE",
        [
            ("The same field is drawn for each role and each tenure band.  ",
             "Where roles disagree, the disagreement is the finding.", ""),
            ("", "", "Teams scoring their own decision autonomy at 2 while managers score it at 4 is a "
                     "diagnosis in itself — and usually more actionable than either number."),
        ],
        [],
    ),
    (
        "WHAT IT NEVER ASKS",
        [
            ("No scored item names a ", "framework, ceremony, role or tool", ".  Not one."),
            ("Every item describes a function — the job being done.", "", ""),
            ("", "Function is what we cannot lose without noticing. Form is what is on the table.", ""),
        ],
        [],
    ),
]

TITLE_BLOCK = [
    ("INSTRUMENT", "Coordination Model Diagnostic — anonymous, ~8–10 min, analysed by role"),
    ("ITEMS", "12 coordination functions × 2 independent ratings"),
    ("OUTPUT", "This field, populated · role-divergence table · ranked keep / fix / cut list"),
    ("STATUS", "Blank — no option decided in advance"),
]


def slide_reading(prs):
    slide = new_slide(prs)
    head(
        slide,
        "Reading the field",
        "Three things the matrix alone does not say, and the record of what the instrument is.",
        ["COORDINATION MODEL DIAGNOSTIC", "HOW TO READ IT"],
    )

    gap = Inches(0.52)
    col_w = Emu(int((W - 2 * M - 2 * gap) / 3))
    top = Inches(1.94)

    for i, (title, paras, bullets) in enumerate(NOTES):
        x = Emu(M + i * (col_w + gap))
        rect(slide, x, top, Inches(0.46), Pt(3.5), fill=INK)

        tf = textbox(slide, x, Emu(top + Inches(0.24)), col_w, Inches(0.26))
        write(tf, title, font=DISPLAY, size=13, bold=True, color=INK, spacing=170, first=True)

        tf = textbox(slide, x, Emu(top + Inches(0.62)), col_w, Inches(1.9))
        for j, (lead, emph, tail) in enumerate(paras):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.line_spacing = 1.3
            para.space_before = Pt(0 if j == 0 else 7)
            for chunk, is_emph in ((lead, False), (emph, True), (tail, False)):
                if not chunk:
                    continue
                run = para.add_run()
                run.text = chunk
                run.font.name = BODY
                run.font.size = Pt(13)
                run.font.bold = is_emph
                run.font.color.rgb = INK if is_emph else INK2

        for k, (swatch, text) in enumerate(bullets):
            y = Emu(top + Inches(1.68) + k * Inches(0.46))
            rect(slide, x, Emu(y + Inches(0.055)), Inches(0.12), Inches(0.12), fill=swatch)
            tf = textbox(slide, Emu(x + Inches(0.24)), y, Emu(col_w - Inches(0.24)), Inches(0.42))
            write(tf, text, size=13, color=INK2, line=1.25, first=True)

    # title block, bottom-right where a drawing's title block belongs
    tb_w, row_h, key_w = Inches(7.4), Inches(0.33), Inches(1.32)
    tb_x, tb_y = Emu(W - M - tb_w), Inches(5.28)
    tb_h = Emu(row_h * len(TITLE_BLOCK))

    for i, (key, val) in enumerate(TITLE_BLOCK):
        y = Emu(tb_y + i * row_h)
        if i:
            hairline(slide, tb_x, y, tb_w, RULE, 0.75)
        tf = textbox(slide, Emu(tb_x + Inches(0.12)), Emu(y + Inches(0.085)), key_w, Inches(0.22))
        write(tf, key, font=MONO, size=9, color=INK3, spacing=60, first=True)
        tf = textbox(slide, Emu(tb_x + key_w + Inches(0.12)), Emu(y + Inches(0.085)),
                     Emu(tb_w - key_w - Inches(0.24)), Inches(0.24))
        write(tf, val, font=MONO, size=9.5, color=INK2, first=True)

    rect(slide, Emu(tb_x + key_w), tb_y, Pt(0.75), tb_h, fill=RULE)
    rect(slide, tb_x, tb_y, tb_w, tb_h, fill=None, line_color=INK, line_pt=1.0)

    footer(slide, 2)
    slide.notes_slide.notes_text_frame.text = (
        "Cost overlay: low importance AND costly is the safe early win. High importance AND costly means "
        "fix it, don't kill it.\n\n"
        "Role divergence is not noise — it is the finding, and usually the most actionable one.\n\n"
        "Neutrality: not one scored item names a framework, ceremony, role or tool. Invite them to check it. "
        "Function is what we cannot lose without noticing; form is what is on the table."
    )
    return slide


# ---------------------------------------------------------------- main


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_field(prs)
    slide_reading(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
